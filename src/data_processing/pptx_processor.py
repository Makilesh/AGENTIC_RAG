"""
PPTX Processor for Agentic RAG System.
"""

import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt

from ..utils import get_data_processing_logger
from .semantic_chunker import Chunk, SemanticChunker

logger = get_data_processing_logger()


@dataclass
class PPTXSlide:
    """
    Represents a single slide from a PowerPoint presentation.
    
    Attributes:
        slide_number: 1-indexed slide number.
        title: Slide title text.
        content: Body content of the slide.
        speaker_notes: Speaker notes associated with the slide.
        has_images: Whether the slide contains images.
        has_tables: Whether the slide contains tables.
    """
    slide_number: int
    title: str
    content: str
    speaker_notes: str = ""
    has_images: bool = False
    has_tables: bool = False


@dataclass
class PPTXDocument:
    """
    Represents a processed PowerPoint presentation.
    
    Attributes:
        file_name: Name of the PPTX file.
        file_path: Full path to the file.
        total_slides: Number of slides in the presentation.
        slides: List of PPTXSlide objects.
        full_text: Combined text from all slides.
        metadata: Presentation-level metadata.
    """
    file_name: str
    file_path: str
    total_slides: int
    slides: List[PPTXSlide]
    full_text: str
    metadata: dict = field(default_factory=dict)


class PPTXProcessor:
    """
    Process PowerPoint presentations for RAG ingestion.
    
    Extracts slide content, titles, and speaker notes with
    slide-level context preservation.
    """
    
    def __init__(self, chunker: Optional[SemanticChunker] = None):
        """
        Initialize the PPTX processor.
        
        Args:
            chunker: SemanticChunker instance. If None, creates a new one.
        """
        self.chunker = chunker or SemanticChunker()
        logger.info("PPTXProcessor initialized")
    
    def _extract_shape_text(self, shape) -> str:
        """
        Extract text from a PowerPoint shape.
        
        Args:
            shape: python-pptx Shape object.
            
        Returns:
            Extracted text content.
        """
        text_parts = []
        
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                para_text = ""
                for run in paragraph.runs:
                    para_text += run.text
                if para_text.strip():
                    text_parts.append(para_text.strip())
        
        return "\n".join(text_parts)
    
    def _extract_table_text(self, shape) -> str:
        """
        Extract text from a PowerPoint table.
        
        Args:
            shape: python-pptx Table shape.
            
        Returns:
            Table content as markdown.
        """
        if not shape.has_table:
            return ""
        
        table = shape.table
        rows = []
        
        for row in table.rows:
            cells = []
            for cell in row.cells:
                cell_text = cell.text.replace("|", "\\|").strip()
                cell_text = re.sub(r'\s+', ' ', cell_text)
                cells.append(cell_text)
            rows.append(cells)
        
        if not rows:
            return ""
        
        # Build markdown table
        lines = []
        
        # Header row
        header = "| " + " | ".join(rows[0]) + " |"
        lines.append(header)
        
        # Separator
        separator = "|" + "|".join(["---"] * len(rows[0])) + "|"
        lines.append(separator)
        
        # Data rows
        for row in rows[1:]:
            while len(row) < len(rows[0]):
                row.append("")
            row_str = "| " + " | ".join(row[:len(rows[0])]) + " |"
            lines.append(row_str)
        
        return "\n".join(lines)
    
    def _extract_speaker_notes(self, slide) -> str:
        """
        Extract speaker notes from a slide.
        
        Args:
            slide: python-pptx Slide object.
            
        Returns:
            Speaker notes text.
        """
        if not slide.has_notes_slide:
            return ""
        
        notes_slide = slide.notes_slide
        notes_frame = notes_slide.notes_text_frame
        
        if notes_frame:
            return notes_frame.text.strip()
        
        return ""
    
    def _get_slide_title(self, slide) -> str:
        """
        Extract the title from a slide.
        
        Args:
            slide: python-pptx Slide object.
            
        Returns:
            Slide title text.
        """
        if slide.shapes.title:
            return slide.shapes.title.text.strip()
        
        # Fallback: look for first text shape
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                return shape.text.strip()[:100]
        
        return f"Slide"
    
    def extract(self, file_path: str) -> PPTXDocument:
        """
        Extract content from a PPTX file.
        
        Args:
            file_path: Path to the PPTX file.
            
        Returns:
            PPTXDocument containing extracted content.
        """
        path = Path(file_path)
        
        if not path.exists():
            raise FileNotFoundError(f"PPTX file not found: {file_path}")
        
        if not path.suffix.lower() == '.pptx':
            raise ValueError(f"Not a PPTX file: {file_path}")
        
        logger.info(f"Processing PPTX: {path.name}")
        
        prs = Presentation(file_path)
        slides = []
        full_text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            title = self._get_slide_title(slide)
            content_parts = []
            has_images = False
            has_tables = False
            
            for shape in slide.shapes:
                # Check for images
                if shape.shape_type == 13:  # PICTURE
                    has_images = True
                    continue
                
                # Check for tables
                if shape.has_table:
                    has_tables = True
                    table_text = self._extract_table_text(shape)
                    if table_text:
                        content_parts.append(f"\n[TABLE]\n{table_text}\n[/TABLE]\n")
                    continue
                
                # Extract regular text
                text = self._extract_shape_text(shape)
                if text and text != title:
                    content_parts.append(text)
            
            # Get speaker notes
            speaker_notes = self._extract_speaker_notes(slide)
            
            # Combine content
            content = "\n".join(content_parts)
            
            pptx_slide = PPTXSlide(
                slide_number=slide_num,
                title=title,
                content=content,
                speaker_notes=speaker_notes,
                has_images=has_images,
                has_tables=has_tables
            )
            slides.append(pptx_slide)
            
            # Build full text representation
            slide_text = f"Slide {slide_num}: {title}\n{content}"
            if speaker_notes:
                slide_text += f"\n\n[Speaker Notes: {speaker_notes}]"
            full_text_parts.append(slide_text)
        
        # Combine full text
        full_text = "\n\n---\n\n".join(full_text_parts)
        
        # Build metadata
        metadata = {
            "source_type": "pptx",
            "file_name": path.name,
            "file_path": str(path.absolute()),
            "total_slides": len(slides),
            "has_tables": any(s.has_tables for s in slides),
            "has_images": any(s.has_images for s in slides),
            "has_notes": any(s.speaker_notes for s in slides)
        }
        
        logger.info(
            f"Extracted {len(slides)} slides from {path.name}",
        )
        
        return PPTXDocument(
            file_name=path.name,
            file_path=str(path.absolute()),
            total_slides=len(slides),
            slides=slides,
            full_text=full_text,
            metadata=metadata
        )
    
    def process(self, file_path: str) -> List[Chunk]:
        """
        Extract and chunk a PPTX presentation.
        
        Each slide becomes a chunk boundary to maintain context.
        
        Args:
            file_path: Path to the PPTX file.
            
        Returns:
            List of Chunk objects ready for embedding.
        """
        # Extract content
        pptx_doc = self.extract(file_path)
        
        # Prepare base metadata
        base_metadata = {
            "source_type": "pptx",
            "file_name": pptx_doc.file_name,
            "file_path": pptx_doc.file_path,
            "total_slides": pptx_doc.total_slides
        }
        
        all_chunks = []
        
        # Process slide by slide
        for slide in pptx_doc.slides:
            slide_metadata = {
                **base_metadata,
                "slide_number": slide.slide_number,
                "slide_title": slide.title,
                "has_notes": bool(slide.speaker_notes),
                "has_tables": slide.has_tables,
                "has_images": slide.has_images
            }
            
            # Build slide text
            slide_text = f"Slide {slide.slide_number}: {slide.title}\n\n{slide.content}"
            
            # Add speaker notes if present
            if slide.speaker_notes:
                slide_text += f"\n\nSpeaker Notes: {slide.speaker_notes}"
            
            if not slide_text.strip():
                continue
            
            # Each slide as a chunk (slides are natural boundaries)
            chunk = Chunk(
                text=slide_text,
                chunk_id=f"{pptx_doc.file_name}_slide{slide.slide_number}",
                chunk_index=len(all_chunks),
                metadata=slide_metadata
            )
            all_chunks.append(chunk)
            
            # If slide content is very long, apply semantic chunking
            if len(slide_text.split()) > 500:
                sub_chunks = self.chunker.chunk(
                    slide_text,
                    f"{pptx_doc.file_name}_slide{slide.slide_number}",
                    slide_metadata
                )
                all_chunks.pop()  # Remove the single slide chunk
                all_chunks.extend(sub_chunks)
        
        # Re-index all chunks
        for i, chunk in enumerate(all_chunks):
            chunk.chunk_index = i
            chunk.metadata["chunk_index"] = i
            chunk.metadata["total_chunks"] = len(all_chunks)
        
        logger.info(
            f"Created {len(all_chunks)} chunks from PPTX {pptx_doc.file_name}",
        )
        
        return all_chunks
